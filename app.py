import streamlit as st
import os
import io
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from query_engine import get_query_engine, format_sources

load_dotenv()

try:
    if hasattr(st, "secrets") and len(st.secrets) > 0:
        for key, value in st.secrets.items():
            os.environ[key] = value
except Exception:
    pass

st.set_page_config(
    page_title="ACENOS Knowledge Agent",
    page_icon="🏦",
    layout="wide"
)

st.title("🏦 ACENOS Knowledge Agent")
st.caption("Interroge la base de connaissances ACENOS · Sources citées automatiquement")


@st.cache_resource
def load_engine():
    return get_query_engine()


engine = load_engine()

# --- Sidebar ---
with st.sidebar:
    st.header("⚙️ Paramètres")

    format_sortie = st.selectbox(
        "Format de réponse",
        ["Texte structuré", "Présentation PPT", "Mail professionnel"],
    )

    st.divider()

    with st.expander("✏️ Modifier le prompt système", expanded=False):
        prompt_path = os.path.join(os.path.dirname(__file__), "prompt.txt")
        try:
            with open(prompt_path, "r", encoding="utf-8") as f:
                prompt_actuel = f.read()
        except FileNotFoundError:
            prompt_actuel = "Fichier prompt.txt introuvable."

        nouveau_prompt = st.text_area(
            "Prompt système",
            value=prompt_actuel,
            height=350,
        )
        if st.button("✅ Appliquer le prompt"):
            with open(prompt_path, "w", encoding="utf-8") as f:
                f.write(nouveau_prompt)
            st.cache_resource.clear()
            st.success("Prompt mis à jour — agent rechargé.")
            st.rerun()

    st.divider()

    if st.button("🗑️ Effacer la conversation"):
        st.session_state.messages = []
        st.rerun()

    st.caption("ACENOS Consulting · Knowledge Agent v1.0")

# --- Historique ---
if "messages" not in st.session_state:
    st.session_state.messages = []

if not st.session_state.messages:
    with st.chat_message("assistant"):
        st.markdown(
            "Bonjour 👋 Je suis l'Agent de Connaissances ACENOS.\n\n"
            "Pose-moi une question sur nos missions, formations, retours d'expérience "
            "ou notre veille réglementaire. Je cite toujours mes sources.\n\n"
            "**Exemples :**\n"
            "- *Quelles sont nos meilleures pratiques pour une due diligence ESG en banque ?*\n"
            "- *Génère les slides d'un pitch pour une mission de transformation comptable*\n"
            "- *Rédige un mail de présentation de nos offres DORA à un DAF bancaire*"
        )

for msg in st.session_state.messages:
    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])


def generer_ppt(question: str, contenu: str) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_titre = prs.slides.add_slide(prs.slide_layouts[0])
    slide_titre.shapes.title.text = question[:80]
    slide_titre.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
    slide_titre.placeholders[1].text = "Base de connaissances ACENOS · Knowledge Agent"

    blocs = [b.strip() for b in contenu.split("\n\n") if b.strip()]
    for bloc in blocs[:8]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        lignes = bloc.split("\n")
        titre = lignes[0].lstrip("#").strip()[:80]
        slide.shapes.title.text = titre
        slide.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
        corps = "\n".join(lignes[1:]).strip()
        if corps and len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.text = corps[:800]
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(16)

    slide_fin = prs.slides.add_slide(prs.slide_layouts[0])
    slide_fin.shapes.title.text = "Sources & Contacts"
    slide_fin.placeholders[1].text = "Document généré par ACENOS Knowledge Agent\ncontact@acenos.fr · www.acenos.fr"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --- Zone de saisie ---
if question := st.chat_input("Pose ta question à la base de connaissances ACENOS..."):

    st.session_state.messages.append({"role": "user", "content": question})
    with st.chat_message("user"):
        st.markdown(question)

    if format_sortie == "Présentation PPT":
        prompt = (
            f"{question}\n\n"
            "Structure ta réponse comme un plan PowerPoint :\n"
            "- Un titre de slide par grande idée (précédé de ##)\n"
            "- 3 à 5 bullet points concis par slide\n"
            "- Maximum 8 slides\n"
            "- Commence par 'Contexte & Enjeux', termine par 'Recommandations'"
        )
    elif format_sortie == "Mail professionnel":
        prompt = (
            f"{question}\n\n"
            "Rédige ta réponse sous forme d'un mail professionnel :\n"
            "- Objet : [objet du mail]\n"
            "- Corps structuré avec introduction, développement et conclusion\n"
            "- Signature : [Prénom Nom] | ACENOS Consulting | contact@acenos.fr"
        )
    else:
        prompt = question

    with st.chat_message("assistant"):
        with st.spinner("Recherche dans la base de connaissances ACENOS..."):
            try:
                response = engine.query(prompt)
                contenu_reponse = str(response)

                st.markdown(contenu_reponse)

                with st.expander("📎 Sources utilisées", expanded=True):
                    st.markdown(format_sources(response))

                if format_sortie == "Présentation PPT":
                    ppt_bytes = generer_ppt(question, contenu_reponse)
                    st.download_button(
                        label="⬇️ Télécharger la présentation PowerPoint",
                        data=ppt_bytes,
                        file_name="acenos_presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )

                st.session_state.messages.append({
                    "role": "assistant",
                    "content": contenu_reponse
                })

            except Exception as e:
                st.error(f"Erreur lors de la requête : {str(e)}")
                st.info("Vérifie que tes clés API sont configurées et que la base est indexée.")